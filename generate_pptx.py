import re
from pptx import Presentation
from pptx.util import Inches, Pt

SRC = "slides.md"
OUT = "slides.pptx"

with open(SRC, "r", encoding="utf-8") as f:
    text = f.read()

# Split slides by markdown slide separator
parts = re.split(r"\n---\n", text)
parts = [p.strip() for p in parts if p.strip()]

prs = Presentation()

for i, part in enumerate(parts):
    lines = [l.rstrip() for l in part.splitlines()]
    # find title (first non-empty line)
    title = ""
    body_lines = []
    notes_lines = []
    mode = "body"
    for ln in lines:
        if not ln.strip():
            continue
        if not title:
            # remove leading Markdown heading marks
            title = re.sub(r"^#{1,6}\s*", "", ln).strip()
            continue
        if ln.strip().startswith("Notes:"):
            mode = "notes"
            # if there is text after 'Notes:' on same line
            after = ln.split("Notes:", 1)[1].strip()
            if after:
                notes_lines.append(after)
            continue
        if mode == "body":
            body_lines.append(ln)
        else:
            notes_lines.append(ln)

    # choose layout: first slide title layout, others title+content
    layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    # set title
    if title:
        try:
            slide.shapes.title.text = title
        except Exception:
            pass

    # add body content
    if body_lines:
        # find content placeholder
        try:
            body = slide.shapes.placeholders[1].text_frame
            # clear default
            body.clear()
        except Exception:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
            body = txBox.text_frame
        for idx, bl in enumerate(body_lines):
            p = body.add_paragraph() if idx != 0 else body.paragraphs[0]
            p.text = bl
            p.level = 0
            p.font.size = Pt(18)

    # add speaker notes
    if notes_lines:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = "\n".join(notes_lines)

prs.save(OUT)
print(f"Saved {OUT}")
